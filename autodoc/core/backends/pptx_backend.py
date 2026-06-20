"""PPTX backend.

A paragraph (inside a shape text frame) whose text matches the question
pattern is a question. We replace a blank in place or append the answer
to the paragraph, copying run formatting. python-pptx round-trips the
presentation, so slide layouts and styles are preserved.
"""
from __future__ import annotations

import json
import re
from typing import Dict, List

from pptx import Presentation

from ..models import Anchor, Answer, DocFormat, EditResult, Question, QuestionType

_QUESTION_RE = re.compile(r"^\s*(\d+)\s*[.、)）]\s*(\S.*)$")
_BLANK_RE = re.compile(r"[_＿]{2,}")


def _para_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def parse_questions(path: str) -> List[Question]:
    prs = Presentation(path)
    questions: List[Question] = []
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                text = _para_text(paragraph)
                m = _QUESTION_RE.match(text)
                if not m:
                    continue
                number = m.group(1)
                blank = _BLANK_RE.search(text)
                mode = "replace_blank" if blank else "append"
                qtype = (
                    QuestionType.fill_blank if blank else QuestionType.short_answer
                )
                locator = json.dumps(
                    {
                        "slide": s_idx,
                        "shape": sh_idx,
                        "para": p_idx,
                        "mode": mode,
                    }
                )
                questions.append(
                    Question(
                        id=f"q{number}",
                        number=number,
                        text=text.strip(),
                        qtype=qtype,
                        anchor=Anchor(
                            format=DocFormat.pptx, locator=locator,
                            detail={"mode": mode},
                        ),
                    )
                )
    return questions


def _replace_blank_in_para(paragraph, answer: str) -> None:
    pos = 0
    for run in paragraph.runs:
        s, e = pos, pos + len(run.text)
        blank = _BLANK_RE.search(run.text)
        if blank:
            run.text = _BLANK_RE.sub(answer, run.text, count=1)
            return
        pos = e
    # blank spanned runs or not found: append
    if paragraph.runs:
        paragraph.runs[-1].text += " " + answer


def _append_to_para(paragraph, answer: str) -> None:
    if paragraph.runs:
        last = paragraph.runs[-1]
        sep = "" if last.text.endswith((" ", "：", ":")) else " "
        last.text = last.text + sep + answer


def write_answers(
    path: str, answers: List[Answer], questions: List[Question], output: str = None
) -> EditResult:
    prs = Presentation(path)
    slides = list(prs.slides)
    by_id: Dict[str, Answer] = {a.question_id: a for a in answers}
    written = 0
    for q in questions:
        ans = by_id.get(q.id)
        if ans is None:
            continue
        loc = json.loads(q.anchor.locator)
        shape = list(slides[loc["slide"]].shapes)[loc["shape"]]
        paragraph = shape.text_frame.paragraphs[loc["para"]]
        if loc["mode"] == "replace_blank":
            _replace_blank_in_para(paragraph, ans.text)
        else:
            _append_to_para(paragraph, ans.text)
        written += 1
    target = output or path
    prs.save(target)
    return EditResult(
        path=target, format=DocFormat.pptx, answers_written=written,
        in_place=(target == path),
    )
