"""Phase 4 tests: txt/md, xlsx, pptx, pdf parsing + in-place editing."""
from __future__ import annotations

import shutil
from pathlib import Path

import pytest

from autodoc.core.edit import write_answers
from autodoc.core.models import Answer, DocFormat
from autodoc.core.parse import parse_document


# ---------- text / markdown ----------

def test_txt_parse_and_edit_preserves_other_lines(sample_txt: Path):
    questions = parse_document(str(sample_txt)).questions
    assert [q.number for q in questions] == ["1", "2"]
    answers = [Answer(question_id="q1", text="北京"), Answer(question_id="q2", text="H2O")]
    result = write_answers(str(sample_txt), answers, questions)
    assert result.in_place
    text = sample_txt.read_text(encoding="utf-8")
    assert "1. 中国的首都是北京。" in text
    assert "2. 请简述水的化学式。 H2O" in text
    # untouched lines remain exactly
    assert "单元测验\n" in text
    assert "说明：请作答。\n" in text
    assert "____" not in text


def test_md_parse_and_edit(sample_md: Path):
    questions = parse_document(str(sample_md)).questions
    assert [q.number for q in questions] == ["1", "2"]
    answers = [Answer(question_id="q1", text="月球"), Answer(question_id="q2", text="将光能转化为化学能")]
    write_answers(str(sample_md), answers, questions)
    text = sample_md.read_text(encoding="utf-8")
    assert "# 测验" in text  # heading preserved
    assert "1. 地球的卫星是月球。" in text


# ---------- xlsx ----------

def test_xlsx_parse_and_edit(sample_xlsx: Path):
    import openpyxl

    questions = parse_document(str(sample_xlsx)).questions
    assert [q.number for q in questions] == ["1", "2"]
    answers = [Answer(question_id="q1", text="北京"), Answer(question_id="q2", text="H2O")]
    result = write_answers(str(sample_xlsx), answers, questions)
    assert result.in_place

    wb = openpyxl.load_workbook(str(sample_xlsx))
    ws = wb.active
    assert ws["A2"].value == "1. 中国的首都是北京。"   # blank replaced in place
    assert ws["B3"].value == "H2O"                     # short answer to right cell
    # preserved: header text, bold font, column width
    assert ws["A1"].value == "题目"
    assert ws["A1"].font.bold is True
    assert ws.column_dimensions["A"].width == 40
    wb.close()


# ---------- pptx ----------

def test_pptx_parse_and_edit(sample_pptx: Path):
    from pptx import Presentation

    questions = parse_document(str(sample_pptx)).questions
    assert [q.number for q in questions] == ["1", "2"]
    answers = [Answer(question_id="q1", text="北京"), Answer(question_id="q2", text="H2O")]
    write_answers(str(sample_pptx), answers, questions)

    prs = Presentation(str(sample_pptx))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append("".join(r.text for r in p.runs))
    joined = "\n".join(texts)
    assert "1. 中国的首都是北京。" in joined
    assert "H2O" in joined


# ---------- pdf form ----------

def test_pdf_form_fill(sample_pdf_form: Path):
    import fitz

    questions = parse_document(str(sample_pdf_form)).questions
    assert any(q.id == "q1" for q in questions)
    answers = [Answer(question_id="q1", text="北京")]
    result = write_answers(str(sample_pdf_form), answers, questions)
    assert result.answers_written == 1

    doc = fitz.open(str(sample_pdf_form))
    values = []
    pages = doc.page_count
    for pno in range(doc.page_count):
        for w in doc[pno].widgets() or []:
            values.append(w.field_value)
    doc.close()
    assert "北京" in values
    assert pages == 1  # page count unchanged


# ---------- pdf text ----------

def test_pdf_text_insert_preserves_pages(sample_pdf_text: Path, tmp_path: Path):
    import fitz

    original = tmp_path / "orig.pdf"
    shutil.copyfile(str(sample_pdf_text), str(original))

    questions = parse_document(str(sample_pdf_text)).questions
    assert [q.number for q in questions] == ["1", "2"]
    answers = [Answer(question_id="q1", text="北京"), Answer(question_id="q2", text="H2O")]
    result = write_answers(str(sample_pdf_text), answers, questions)
    assert result.answers_written >= 1

    before = fitz.open(str(original))
    after = fitz.open(str(sample_pdf_text))
    try:
        assert before.page_count == after.page_count  # page count unchanged
        assert before[0].rect == after[0].rect        # page geometry unchanged
        after_text = after[0].get_text("text")
        # original question text still present (we only added content)
        assert "请简述水的化学式" in after_text
        assert "北京" in after_text
    finally:
        before.close()
        after.close()


def test_format_detection_round_trip():
    from autodoc.core.formats import detect_format

    assert detect_format("a.docx") == DocFormat.docx
    assert detect_format("a.PDF") == DocFormat.pdf
    assert detect_format("a.md") == DocFormat.md
    with pytest.raises(ValueError):
        detect_format("a.rtf")
