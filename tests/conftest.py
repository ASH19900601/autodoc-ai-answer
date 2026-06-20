"""Shared pytest fixtures: programmatically built sample documents."""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

import docx
import pytest
from docx.shared import Pt


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """A styled .docx with mixed question types, blanks and formatting."""
    document = docx.Document()

    # custom section margins to verify section properties survive editing
    section = document.sections[0]
    section.left_margin = Pt(72)
    section.right_margin = Pt(72)

    document.add_heading("单元测验", level=1)

    document.add_paragraph("说明：请在横线处作答。")  # non-question

    document.add_paragraph("1. 中国的首都是____。")  # fill_blank
    document.add_paragraph("2. 请简述水的化学式。")  # short_answer

    p3 = document.add_paragraph("3. 地球唯一的天然卫星是")
    run = p3.add_run("____")
    run.bold = True
    run.font.size = Pt(14)
    p3.add_run("。")

    path = tmp_path / "quiz.docx"
    document.save(str(path))
    return path


def read_zip_member(path, member: str) -> bytes:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        if member not in names:
            return b""
        return zf.read(member)


def zip_members(path):
    with zipfile.ZipFile(path) as zf:
        return set(zf.namelist())


def copy_file(src, dst) -> Path:
    shutil.copyfile(str(src), str(dst))
    return Path(dst)


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    path = tmp_path / "quiz.txt"
    content = (
        "单元测验\n"
        "说明：请作答。\n"
        "1. 中国的首都是____。\n"
        "2. 请简述水的化学式。\n"
    )
    path.write_text(content, encoding="utf-8")
    return path


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    path = tmp_path / "quiz.md"
    path.write_text(
        "# 测验\n\n1. 地球的卫星是____。\n2. 简述光合作用。\n", encoding="utf-8"
    )
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path):
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "题目"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "1. 中国的首都是____。"
    ws["A3"] = "2. 请简述水的化学式。"
    ws.column_dimensions["A"].width = 40
    path = tmp_path / "quiz.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = box.text_frame
    tf.paragraphs[0].text = "1. 中国的首都是____。"
    p2 = tf.add_paragraph()
    p2.text = "2. 请简述水的化学式。"
    path = tmp_path / "quiz.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_pdf_form(tmp_path: Path):
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(fitz.Point(72, 72), "1. 中国的首都是：")
    widget = fitz.Widget()
    widget.field_name = "q1"
    widget.field_type = fitz.PDF_WIDGET_TYPE_TEXT
    widget.rect = fitz.Rect(200, 60, 400, 80)
    page.add_widget(widget)
    path = tmp_path / "form.pdf"
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def sample_pdf_text(tmp_path: Path):
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(fitz.Point(72, 100), "1. 中国的首都是 ____ 。", fontsize=12, fontname="china-s")
    page.insert_text(fitz.Point(72, 140), "2. 请简述水的化学式。", fontsize=12, fontname="china-s")
    path = tmp_path / "text.pdf"
    doc.save(str(path))
    doc.close()
    return path
